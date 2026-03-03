import os
import io
import json
import math
import time
import base64
import asyncio
import datetime
from typing import List, Optional, Dict
from functools import partial
from threading import Lock

import joblib
import pandas as pd
import numpy as np
import faiss
import cv2
from PIL import Image
import torch
import mediapipe as mp
from mediapipe.tasks.python import vision as mp_vision
from mediapipe.tasks import python as mp_tasks

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import train_test_split
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from ultralytics import YOLO

from dotenv import load_dotenv

# ✅ New Gemini SDK (replaces deprecated google.generativeai)
from google import genai

# ----------------------------
# Env
# ----------------------------
load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-1.5-flash")  # start with flash (fast + usually available)

# Create Gemini client only if key exists (so server can still run without it)
gemini_client = genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None

app = FastAPI(title="Campus AI Analytics & RAG API")

# CORS - allow frontend to call this service directly
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ----------------------------
# GPU / Device Setup
# ----------------------------
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
print(f"Using device: {DEVICE}")
if DEVICE == "cuda":
    print(f"GPU: {torch.cuda.get_device_name(0)}")


@app.get("/")
async def root():
    return {
        "service": "Campus AI Analytics & RAG API",
        "status": "running",
        "device": DEVICE,
        "endpoints": {
            "ml": ["/train", "/predict"],
            "rag": ["/process-material", "/chat"],
            "proctoring": [
                "/api/proctor/detect-objects",
                "/api/proctor/head-pose",
                "/api/proctor/analyze",
                "/api/proctor/health",
            ],
        },
    }

# ----------------------------
# YOLOv8m Model for Object Detection (Proctoring)
# ----------------------------
print("Loading YOLOv8m model for proctoring (balanced speed/accuracy)...")
yolo_model = YOLO("yolov8m.pt")  # medium model - best balance for MX150

# COCO class IDs for suspicious objects during exams
SUSPICIOUS_OBJECTS = {
    67: "cell phone",
    73: "book",
    63: "laptop",
    65: "remote",
    74: "clock",
    0: "person",  # tracked separately for multiple-person detection
}

# ----------------------------
# MediaPipe Face Detection + Face Landmarker (Tasks API)
# ----------------------------
print("Loading MediaPipe Tasks API models...")
TASK_MODELS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "task_models")

face_detector = mp_vision.FaceDetector.create_from_options(
    mp_vision.FaceDetectorOptions(
        base_options=mp_tasks.BaseOptions(
            model_asset_path=os.path.join(TASK_MODELS_DIR, "blaze_face_short_range.tflite")
        ),
        min_detection_confidence=0.7,
        running_mode=mp_vision.RunningMode.IMAGE,
    )
)

face_landmarker = mp_vision.FaceLandmarker.create_from_options(
    mp_vision.FaceLandmarkerOptions(
        base_options=mp_tasks.BaseOptions(
            model_asset_path=os.path.join(TASK_MODELS_DIR, "face_landmarker.task")
        ),
        output_face_blendshapes=True,
        num_faces=2,
        min_face_detection_confidence=0.7,
        min_face_presence_confidence=0.7,
        running_mode=mp_vision.RunningMode.IMAGE,
    )
)
print("MediaPipe Tasks API models loaded.")

# ----------------------------
# Embedding Model
# ----------------------------
print("Loading Embedding Model...")
embed_model = SentenceTransformer("all-MiniLM-L6-v2")

MODEL_DIR = "models"
INDEX_DIR = "indices"
os.makedirs(MODEL_DIR, exist_ok=True)
os.makedirs(INDEX_DIR, exist_ok=True)

METADATA_PATH = os.path.join(MODEL_DIR, "metadata.json")

# FAISS Indices storage (in-memory)
# courseId -> {"index": faiss.IndexFlatL2(dim), "chunks": [chunk_info, ...]}
indices: Dict[str, Dict] = {}


# ----------------------------
# Temporal Violation Tracking (Per-Attempt)
# ----------------------------
class ViolationBuffer:
    """Thread-safe in-memory store tracking temporal violation patterns per quiz attempt."""

    def __init__(self):
        self._lock = Lock()
        self._buffers: Dict[str, Dict] = {}

    def get_or_create(self, attempt_id: str) -> Dict:
        with self._lock:
            if attempt_id not in self._buffers:
                self._buffers[attempt_id] = {
                    "no_face_start": None,
                    "looking_away_start": None,
                    "looking_away_direction": None,
                    "last_seen": time.time(),
                }
            self._buffers[attempt_id]["last_seen"] = time.time()
            return self._buffers[attempt_id]

    def cleanup_stale(self, max_age_seconds: float = 7200):
        """Remove attempt buffers older than 2 hours."""
        with self._lock:
            now = time.time()
            stale = [k for k, v in self._buffers.items() if now - v["last_seen"] > max_age_seconds]
            for k in stale:
                del self._buffers[k]


violation_buffer = ViolationBuffer()


def _analyze_temporal_patterns(
    attempt_id: str,
    face_count: int,
    looking_away: bool,
    looking_direction: str,
) -> List[Dict]:
    """Track sustained violations over time for a quiz attempt."""
    buf = violation_buffer.get_or_create(attempt_id)
    now = time.time()
    temporal_violations = []

    # SUSTAINED_NO_FACE: no face for >5 continuous seconds
    if face_count == 0:
        if buf["no_face_start"] is None:
            buf["no_face_start"] = now
        else:
            duration = now - buf["no_face_start"]
            if duration >= 5.0:
                temporal_violations.append({
                    "type": "SUSTAINED_NO_FACE",
                    "duration_seconds": round(duration, 1),
                    "details": f"No face detected for {duration:.1f}s continuously",
                })
    else:
        buf["no_face_start"] = None

    # SUSTAINED_GAZE_DEVIATION: looking away >3 continuous seconds
    if looking_away and face_count == 1:
        if buf["looking_away_start"] is None:
            buf["looking_away_start"] = now
            buf["looking_away_direction"] = looking_direction
        else:
            duration = now - buf["looking_away_start"]
            if duration >= 3.0:
                temporal_violations.append({
                    "type": "SUSTAINED_GAZE_DEVIATION",
                    "duration_seconds": round(duration, 1),
                    "details": f"Looking {buf['looking_away_direction']} for {duration:.1f}s continuously",
                })
    else:
        buf["looking_away_start"] = None
        buf["looking_away_direction"] = None

    return temporal_violations


# ----------------------------
# Helpers: metadata
# ----------------------------
def save_metadata(metadata: dict) -> None:
    with open(METADATA_PATH, "w") as f:
        json.dump(metadata, f)


def load_metadata() -> dict:
    if os.path.exists(METADATA_PATH):
        with open(METADATA_PATH, "r") as f:
            return json.load(f)
    return {"current_version": "v1", "models": {}}


# ----------------------------
# ML: Risk prediction schemas
# ----------------------------
class StudentData(BaseModel):
    student_id: str
    subject_id: str
    attendance_percentage: Optional[float] = 0
    average_assignment_score: Optional[float] = 0
    quiz_average: Optional[float] = 0
    gpa: Optional[float] = 0
    classes_missed: Optional[int] = 0
    late_submissions: Optional[int] = 0
    quiz_attempts: Optional[int] = 0
    face_violation_count: Optional[int] = 0
    payment_delay_days: Optional[int] = 0
    previous_exam_score: Optional[float] = 0
    exam_result: Optional[int] = None  # 1 FAIL, 0 PASS


class PredictionResponse(BaseModel):
    student_id: str
    risk_score: float
    risk_level: str
    reasons: List[str]
    version: Optional[str] = None


def _sync_train(data: List[StudentData]):
    """Run blocking ML training in a worker thread."""
    df = pd.DataFrame([d.dict() for d in data])

    # Feature Engineering
    df["attendance_risk"] = (df["attendance_percentage"] < 75).astype(int)
    df["low_gpa"] = (df["gpa"] < 2.5).astype(int)

    X = df.drop(columns=["student_id", "subject_id", "exam_result"])
    y = df["exam_result"]

    if y.isnull().any():
        raise HTTPException(status_code=400, detail="Target variable 'exam_result' is missing in some records.")

    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

    model = RandomForestClassifier(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)

    accuracy = model.score(X_test, y_test)

    metadata = load_metadata()
    version_num = len(metadata["models"]) + 1
    version = f"v{version_num}"
    model_filename = f"model_{version}.pkl"
    model_path = os.path.join(MODEL_DIR, model_filename)

    joblib.dump(model, model_path)

    metadata["current_version"] = version
    metadata["models"][version] = {
        "path": model_path,
        "trained_at": datetime.datetime.now().isoformat(),
        "accuracy": float(accuracy),
        "records_used": len(data),
    }
    save_metadata(metadata)

    return {"message": "Model trained successfully", "version": version, "accuracy": float(accuracy), "records_used": len(data)}


@app.post("/train")
async def train(data: List[StudentData]):
    if len(data) < 10:
        raise HTTPException(status_code=400, detail="Not enough data to train. Need at least 10 records.")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_sync_train, data))


def _sync_predict(data: StudentData):
    """Run blocking ML prediction in a worker thread."""
    metadata = load_metadata()
    current_version = metadata.get("current_version", "v1")
    model_info = metadata.get("models", {}).get(current_version)

    if not model_info or not os.path.exists(model_info["path"]):
        res = fallback_predict(data)
        res["version"] = "fallback"
        return res

    model = joblib.load(model_info["path"])

    df = pd.DataFrame([data.dict()])
    df["attendance_risk"] = (df["attendance_percentage"] < 75).astype(int)
    df["low_gpa"] = (df["gpa"] < 2.5).astype(int)

    X = df.drop(columns=["student_id", "subject_id", "exam_result"])

    risk_score = model.predict_proba(X)[0][1]  # prob of FAIL (class 1)

    if risk_score >= 0.7:
        risk_level = "HIGH"
    elif risk_score >= 0.4:
        risk_level = "MEDIUM"
    else:
        risk_level = "LOW"

    reasons = []
    if data.attendance_percentage < 75:
        reasons.append("Low attendance")
    if data.quiz_average < 50:
        reasons.append("Low quiz performance")
    if data.face_violation_count > 3:
        reasons.append("Multiple face violations")
    if data.average_assignment_score < 50:
        reasons.append("Low assignment scores")

    return {
        "student_id": data.student_id,
        "risk_score": float(risk_score),
        "risk_level": risk_level,
        "reasons": reasons,
        "version": current_version,
    }


@app.post("/predict", response_model=PredictionResponse)
async def predict(data: StudentData):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_sync_predict, data))


def fallback_predict(data: StudentData):
    score = 0.0
    reasons = []

    if data.attendance_percentage < 75:
        score += 0.4
        reasons.append("Low attendance")
    if data.quiz_average < 50:
        score += 0.2
        reasons.append("Low quiz performance")
    if data.average_assignment_score < 50:
        score += 0.2
        reasons.append("Low assignment scores")
    if data.face_violation_count > 3:
        score += 0.1
        reasons.append("Multiple face violations")
    if data.payment_delay_days > 15:
        score += 0.1
        reasons.append("Payment delays")

    score = min(score, 1.0)

    if score >= 0.7:
        risk_level = "HIGH"
    elif score >= 0.4:
        risk_level = "MEDIUM"
    else:
        risk_level = "LOW"

    return {"student_id": data.student_id, "risk_score": score, "risk_level": risk_level, "reasons": reasons}


# ----------------------------
# RAG: Schemas
# ----------------------------
class ChatRequest(BaseModel):
    courseId: str
    question: str
    top_k: Optional[int] = 5


class ChunkData(BaseModel):
    content: str
    metadata: Dict


def extract_text_from_file(file_content: bytes, filename: str) -> List[ChunkData]:
    chunks: List[ChunkData] = []
    ext = filename.split(".")[-1].lower()

    if ext == "pdf":
        reader = PdfReader(io.BytesIO(file_content))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(ChunkData(content=text.strip(), metadata={"page": i + 1}))

    elif ext == "docx":
        doc = Document(io.BytesIO(file_content))
        current_text = ""
        for i, para in enumerate(doc.paragraphs):
            current_text += (para.text or "") + "\n"
            if len(current_text) > 1000:
                chunks.append(ChunkData(content=current_text.strip(), metadata={"paragraph_index": i}))
                current_text = ""
        if current_text.strip():
            chunks.append(ChunkData(content=current_text.strip(), metadata={"paragraph_index": len(doc.paragraphs)}))

    elif ext == "pptx":
        prs = Presentation(io.BytesIO(file_content))
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += (shape.text or "") + " "
            if text.strip():
                chunks.append(ChunkData(content=text.strip(), metadata={"slide": i + 1}))

    else:
        # Generic text
        text = file_content.decode("utf-8", errors="ignore")
        for i in range(0, len(text), 1000):
            piece = text[i : i + 1000].strip()
            if piece:
                chunks.append(ChunkData(content=piece, metadata={"offset": i}))

    return chunks


def _sync_process_material(courseId: str, materialId: str, content: bytes, filename: str) -> dict:
    """Run blocking file parsing + embedding in a worker thread."""
    chunks = extract_text_from_file(content, filename)

    if not chunks:
        return {"message": "No text extracted", "chunks": []}

    texts = [c.content for c in chunks]
    embeddings = embed_model.encode(texts)

    dim = embeddings.shape[1]
    if courseId not in indices:
        indices[courseId] = {"index": faiss.IndexFlatL2(dim), "chunks": []}

    start_idx = len(indices[courseId]["chunks"])
    indices[courseId]["index"].add(embeddings.astype("float32"))

    processed_chunks = []
    for i, c in enumerate(chunks):
        chunk_info = {
            "id": f"{materialId}_{start_idx + i}",
            "materialId": materialId,
            "content": c.content,
            "metadata": c.metadata,
            "global_idx": start_idx + i,
        }
        indices[courseId]["chunks"].append(chunk_info)
        processed_chunks.append(chunk_info)

    return {"message": f"Processed {len(chunks)} chunks", "chunks": processed_chunks}


@app.post("/process-material")
async def process_material(courseId: str = Form(...), materialId: str = Form(...), file: UploadFile = File(...)):
    content = await file.read()
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, partial(_sync_process_material, courseId, materialId, content, file.filename)
    )


def build_prompt(question: str, chunks: list) -> str:
    """Prompt that grounds the answer in lecture materials when available,
    but also allows the model to draw on general knowledge."""
    if chunks:
        sources = []
        for i, c in enumerate(chunks, start=1):
            meta = c.get("metadata", {})
            sources.append(f"[S{i}] meta={meta}\n{c['content']}")

        return f"""
You are a helpful university assistant with broad academic knowledge.

LECTURE MATERIAL CONTEXT (use these as your primary source when relevant):
{chr(10).join(sources)}

RULES:
- Prioritise the lecture material above when it is relevant to the question.
- Add citations like [S1], [S2] after statements that come from the lecture material.
- For questions not covered by the lecture material, answer using your general knowledge.
- Always give a clear, student-friendly answer — never refuse to answer.

QUESTION:
{question}
""".strip()
    else:
        return f"""
You are a helpful university assistant with broad academic knowledge.

Answer the following question in a clear, student-friendly way.

QUESTION:
{question}
""".strip()


def _sync_chat(course_id: str, question: str, top_k: int) -> dict:
    """Run all blocking chat operations (embedding, FAISS, Gemini) in a worker thread."""

    # Try to retrieve relevant lecture material chunks (best-effort — never blocks the answer)
    retrieved_chunks = []
    if course_id in indices and indices[course_id]["chunks"]:
        q_emb = embed_model.encode([question])
        D, I = indices[course_id]["index"].search(q_emb.astype("float32"), top_k)
        for idx in I[0]:
            if idx != -1 and idx < len(indices[course_id]["chunks"]):
                retrieved_chunks.append(indices[course_id]["chunks"][idx])

    prompt = build_prompt(question, retrieved_chunks)

    # Call Gemini
    if not gemini_client:
        answer = "Gemini is not configured. Please set GEMINI_API_KEY (or GOOGLE_API_KEY) in your .env and restart."
    else:
        try:
            resp = gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
            )
            answer = resp.text or "I couldn't generate an answer. Please try again."
        except Exception as e:
            print(f"Error calling Gemini: {e}")
            answer = "I'm sorry, I'm having trouble connecting to Gemini right now."

    return {
        "answer": answer,
        "citations": [
            {
                "source": f"S{i+1}",
                "materialId": c["materialId"],
                "metadata": c["metadata"],
                "snippet": c["content"][:150],
            }
            for i, c in enumerate(retrieved_chunks)
        ],
    }


@app.post("/chat")
async def chat(request: ChatRequest):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, partial(_sync_chat, request.courseId, request.question, request.top_k)
    )


# ----------------------------
# Proctoring: Object Detection
# ----------------------------
class ObjectDetectionRequest(BaseModel):
    image: str  # base64-encoded JPEG frame
    confidence_threshold: Optional[float] = 0.5


class DetectedObject(BaseModel):
    label: str
    confidence: float
    bbox: List[float]  # [x1, y1, x2, y2]


class ObjectDetectionResponse(BaseModel):
    suspicious_objects: List[DetectedObject]
    person_count: int
    is_violation: bool
    violation_details: Optional[str] = None


def _sync_detect_objects(image_b64: str, confidence_threshold: float) -> dict:
    """Run blocking YOLO inference in a worker thread."""
    image_data = base64.b64decode(image_b64)
    image = Image.open(io.BytesIO(image_data))
    frame = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)

    results = yolo_model(
        frame,
        device=DEVICE,
        conf=confidence_threshold,
        verbose=False,
        half=(DEVICE == "cuda"),
    )

    suspicious: List[DetectedObject] = []
    person_count = 0

    for r in results:
        for box in r.boxes:
            cls_id = int(box.cls[0])
            conf = float(box.conf[0])
            bbox = box.xyxy[0].tolist()

            if cls_id == 0:
                person_count += 1
            elif cls_id in SUSPICIOUS_OBJECTS:
                suspicious.append(DetectedObject(
                    label=SUSPICIOUS_OBJECTS[cls_id],
                    confidence=round(conf, 3),
                    bbox=bbox,
                ))

    is_violation = len(suspicious) > 0
    violation_details = None
    if suspicious:
        labels = [obj.label for obj in suspicious]
        violation_details = f"Detected: {', '.join(labels)}"

    return ObjectDetectionResponse(
        suspicious_objects=suspicious,
        person_count=person_count,
        is_violation=is_violation,
        violation_details=violation_details,
    )


@app.post("/api/proctor/detect-objects", response_model=ObjectDetectionResponse)
async def detect_objects(request: ObjectDetectionRequest):
    try:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, partial(_sync_detect_objects, request.image, request.confidence_threshold)
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Object detection failed: {str(e)}")


# ----------------------------
# Proctoring: Head Pose Detection (MediaPipe FaceMesh)
# ----------------------------
class HeadPoseRequest(BaseModel):
    image: str  # base64-encoded JPEG frame


class HeadPoseResponse(BaseModel):
    face_count: int
    head_pose: Optional[Dict] = None  # yaw, pitch, roll in degrees
    looking_away: bool
    looking_direction: str  # "center", "left", "right", "up", "down"
    eyes_closed: bool
    is_violation: bool
    violation_type: Optional[str] = None  # "NO_FACE", "MULTIPLE_FACES", "LOOKING_AWAY", "EYES_CLOSED"
    violation_details: Optional[str] = None


def estimate_head_pose(landmarks, img_w: int, img_h: int) -> Dict:
    """Estimate head pose (yaw, pitch, roll) from MediaPipe FaceMesh landmarks."""
    # Key landmark indices for pose estimation
    # Nose tip, chin, left eye corner, right eye corner, left mouth, right mouth
    face_3d = []
    face_2d = []

    key_indices = [1, 33, 61, 199, 263, 291]

    for idx in key_indices:
        lm = landmarks[idx]
        x, y = int(lm.x * img_w), int(lm.y * img_h)
        face_2d.append([x, y])
        face_3d.append([x, y, lm.z * 3000])

    face_2d = np.array(face_2d, dtype=np.float64)
    face_3d = np.array(face_3d, dtype=np.float64)

    # Camera matrix (approximate)
    focal_length = img_w
    cam_matrix = np.array([
        [focal_length, 0, img_w / 2],
        [0, focal_length, img_h / 2],
        [0, 0, 1],
    ], dtype=np.float64)

    dist_coeffs = np.zeros((4, 1), dtype=np.float64)

    _, rot_vec, _ = cv2.solvePnP(face_3d, face_2d, cam_matrix, dist_coeffs)
    rmat, _ = cv2.Rodrigues(rot_vec)

    angles, _, _, _, _, _ = cv2.RQDecomp3x3(rmat)

    yaw = angles[1] * 360    # left/right
    pitch = angles[0] * 360  # up/down
    roll = angles[2] * 360   # tilt

    return {"yaw": round(yaw, 1), "pitch": round(pitch, 1), "roll": round(roll, 1)}


def detect_eyes_closed(landmarks) -> bool:
    """Detect if eyes are closed using Eye Aspect Ratio (EAR)."""
    def ear(eye_indices):
        pts = [(landmarks[i].x, landmarks[i].y) for i in eye_indices]
        # Vertical distances
        v1 = math.dist(pts[1], pts[5])
        v2 = math.dist(pts[2], pts[4])
        # Horizontal distance
        h = math.dist(pts[0], pts[3])
        return (v1 + v2) / (2.0 * h) if h > 0 else 0

    # MediaPipe FaceMesh eye landmark indices
    left_eye = [33, 160, 158, 133, 153, 144]
    right_eye = [362, 385, 387, 263, 373, 380]

    left_ear = ear(left_eye)
    right_ear = ear(right_eye)
    avg_ear = (left_ear + right_ear) / 2.0

    return avg_ear < 0.18  # threshold for closed eyes


def _sync_detect_head_pose(image_b64: str) -> HeadPoseResponse:
    """Run blocking MediaPipe inference in a worker thread."""
    image_data = base64.b64decode(image_b64)
    image = Image.open(io.BytesIO(image_data))
    frame = np.array(image)
    rgb_frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR) if len(frame.shape) == 3 else frame

    rgb = cv2.cvtColor(rgb_frame, cv2.COLOR_BGR2RGB)
    img_h, img_w = rgb.shape[:2]

    mp_image = mp.Image(image_format=mp.ImageFormat.SRGB, data=rgb)

    face_results = face_detector.detect(mp_image)
    face_count = len(face_results.detections)

    if face_count == 0:
        return HeadPoseResponse(
            face_count=0, looking_away=False, looking_direction="none",
            eyes_closed=False, is_violation=True, violation_type="NO_FACE",
            violation_details="No face detected in frame",
        )

    if face_count > 1:
        return HeadPoseResponse(
            face_count=face_count, looking_away=False, looking_direction="center",
            eyes_closed=False, is_violation=True, violation_type="MULTIPLE_FACES",
            violation_details=f"{face_count} faces detected",
        )

    landmarker_result = face_landmarker.detect(mp_image)

    if not landmarker_result.face_landmarks:
        return HeadPoseResponse(
            face_count=1, looking_away=False, looking_direction="center",
            eyes_closed=False, is_violation=False,
        )

    landmarks = landmarker_result.face_landmarks[0]
    head_pose = estimate_head_pose(landmarks, img_w, img_h)
    eyes_closed = detect_eyes_closed(landmarks)

    yaw = head_pose["yaw"]
    pitch = head_pose["pitch"]
    looking_direction = "center"
    looking_away = False

    if yaw < -15:
        looking_direction = "left"
        looking_away = True
    elif yaw > 15:
        looking_direction = "right"
        looking_away = True
    if pitch < -15:
        looking_direction = "down"
        looking_away = True
    elif pitch > 20:
        looking_direction = "up"
        looking_away = True

    is_violation = looking_away or eyes_closed
    violation_type = None
    violation_details = None

    if looking_away:
        violation_type = "LOOKING_AWAY"
        violation_details = f"Looking {looking_direction} (yaw={yaw:.0f}, pitch={pitch:.0f})"
    elif eyes_closed:
        violation_type = "EYES_CLOSED"
        violation_details = "Eyes appear to be closed"

    return HeadPoseResponse(
        face_count=face_count, head_pose=head_pose, looking_away=looking_away,
        looking_direction=looking_direction, eyes_closed=eyes_closed,
        is_violation=is_violation, violation_type=violation_type,
        violation_details=violation_details,
    )


@app.post("/api/proctor/head-pose", response_model=HeadPoseResponse)
async def detect_head_pose(request: HeadPoseRequest):
    try:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, partial(_sync_detect_head_pose, request.image))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Head pose detection failed: {str(e)}")


# ----------------------------
# Proctoring: Combined Analysis (single endpoint for efficiency)
# ----------------------------
class FullProctorRequest(BaseModel):
    image: str  # base64-encoded JPEG frame
    run_object_detection: bool = True  # skip YOLO on some frames for perf
    confidence_threshold: Optional[float] = 0.6
    attempt_id: Optional[str] = None  # enables temporal violation tracking


class FullProctorResponse(BaseModel):
    # Face / Head Pose
    face_count: int
    head_pose: Optional[Dict] = None
    looking_away: bool
    looking_direction: str
    eyes_closed: bool
    # Object Detection
    suspicious_objects: List[DetectedObject]
    person_count: int
    # Violations
    violations: List[Dict]  # [{type, details, weight}]
    total_violation_weight: int
    # Temporal (only when attempt_id is provided)
    temporal_violations: Optional[List[Dict]] = None
    attempt_id: Optional[str] = None


def _sync_full_proctor_analyze(image_b64: str, run_object_detection: bool, confidence_threshold: float, attempt_id: Optional[str] = None) -> FullProctorResponse:
    """Run blocking combined proctoring analysis in a worker thread."""
    image_data = base64.b64decode(image_b64)
    image = Image.open(io.BytesIO(image_data))
    frame_rgb = np.array(image)
    frame_bgr = cv2.cvtColor(frame_rgb, cv2.COLOR_RGB2BGR)
    img_h, img_w = frame_bgr.shape[:2]

    rgb = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2RGB)
    mp_image = mp.Image(image_format=mp.ImageFormat.SRGB, data=rgb)

    violations: List[Dict] = []

    face_results = face_detector.detect(mp_image)
    face_count = len(face_results.detections)

    head_pose = None
    looking_away = False
    looking_direction = "none"
    eyes_closed = False

    if face_count == 0:
        violations.append({"type": "NO_FACE", "details": "No face detected", "weight": 1})
    elif face_count > 1:
        violations.append({"type": "MULTIPLE_FACES", "details": f"{face_count} faces detected", "weight": 2})
    else:
        landmarker_result = face_landmarker.detect(mp_image)
        if landmarker_result.face_landmarks:
            landmarks = landmarker_result.face_landmarks[0]
            head_pose = estimate_head_pose(landmarks, img_w, img_h)
            eyes_closed = detect_eyes_closed(landmarks)

            yaw = head_pose["yaw"]
            pitch = head_pose["pitch"]
            looking_direction = "center"

            if yaw < -15:
                looking_direction = "left"
                looking_away = True
            elif yaw > 15:
                looking_direction = "right"
                looking_away = True
            if pitch < -15:
                looking_direction = "down"
                looking_away = True
            elif pitch > 20:
                looking_direction = "up"
                looking_away = True

            if looking_away:
                violations.append({
                    "type": "LOOKING_AWAY",
                    "details": f"Looking {looking_direction} (yaw={yaw:.0f}, pitch={pitch:.0f})",
                    "weight": 1,
                })

    suspicious: List[DetectedObject] = []
    person_count = 0

    if run_object_detection:
        results = yolo_model(
            frame_bgr,
            device=DEVICE,
            conf=confidence_threshold,
            verbose=False,
            half=(DEVICE == "cuda"),
        )

        for r in results:
            for box in r.boxes:
                cls_id = int(box.cls[0])
                conf = float(box.conf[0])
                bbox = box.xyxy[0].tolist()

                if cls_id == 0:
                    person_count += 1
                elif cls_id in SUSPICIOUS_OBJECTS:
                    label = SUSPICIOUS_OBJECTS[cls_id]
                    suspicious.append(DetectedObject(
                        label=label,
                        confidence=round(conf, 3),
                        bbox=bbox,
                    ))
                    weight = 3 if label == "cell phone" else 2
                    violations.append({
                        "type": "CHEATING_OBJECT",
                        "details": f"Detected: {label} ({conf:.0%})",
                        "weight": weight,
                    })

    # Temporal analysis (only when attempt_id is provided)
    temporal_violations = None
    if attempt_id:
        temporal_violations = _analyze_temporal_patterns(
            attempt_id, face_count, looking_away, looking_direction
        )
        # Add temporal violations to the main violations list
        if temporal_violations:
            for tv in temporal_violations:
                violations.append({
                    "type": tv["type"],
                    "details": tv["details"],
                    "weight": 2,
                })

    total_weight = sum(v["weight"] for v in violations)

    return FullProctorResponse(
        face_count=face_count,
        head_pose=head_pose,
        looking_away=looking_away,
        looking_direction=looking_direction,
        eyes_closed=eyes_closed,
        suspicious_objects=suspicious,
        person_count=person_count,
        violations=violations,
        total_violation_weight=total_weight,
        temporal_violations=temporal_violations,
        attempt_id=attempt_id,
    )


class CleanupBufferRequest(BaseModel):
    attempt_id: Optional[str] = None

@app.post("/api/proctor/cleanup-buffer")
async def cleanup_violation_buffer(request: CleanupBufferRequest = CleanupBufferRequest()):
    """Cleanup violation tracking buffers. If attempt_id is provided, clears only that buffer.
    Otherwise cleans up all stale buffers (>2 hours old)."""
    if request.attempt_id:
        with violation_buffer._lock:
            if request.attempt_id in violation_buffer._buffers:
                del violation_buffer._buffers[request.attempt_id]
        return {"status": "ok", "cleared": request.attempt_id}
    else:
        violation_buffer.cleanup_stale()
        return {"status": "ok"}


@app.post("/api/proctor/analyze", response_model=FullProctorResponse)
async def full_proctor_analyze(request: FullProctorRequest):
    """Combined proctoring analysis - face detection + head pose + optional object detection.
    Call this every frame for face/pose, but set run_object_detection=True only every 5th frame."""
    try:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, partial(_sync_full_proctor_analyze, request.image, request.run_object_detection, request.confidence_threshold, request.attempt_id)
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Proctor analysis failed: {str(e)}")


@app.get("/api/proctor/health")
async def proctor_health():
    return {
        "status": "ok",
        "device": DEVICE,
        "gpu_available": torch.cuda.is_available(),
        "gpu_name": torch.cuda.get_device_name(0) if torch.cuda.is_available() else None,
        "yolo_model": "yolov8m",
        "mediapipe": "tasks-api",
    }


# ----------------------------
# Face Recognition: Identity Verification (DeepFace + ArcFace)
# ----------------------------
FACE_EMBEDDINGS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "face_embeddings")
os.makedirs(FACE_EMBEDDINGS_DIR, exist_ok=True)

ARCFACE_THRESHOLD = 0.68  # Cosine similarity threshold for ArcFace

# Lazy-load DeepFace to avoid slow startup if not needed
_deepface_module = None


def _get_deepface():
    global _deepface_module
    if _deepface_module is None:
        from deepface import DeepFace
        _deepface_module = DeepFace
    return _deepface_module


class FaceEnrollRequest(BaseModel):
    student_id: str
    image: str  # base64-encoded JPEG


class FaceEnrollResponse(BaseModel):
    success: bool
    message: str
    student_id: str


class FaceVerifyRequest(BaseModel):
    student_id: str
    image: str  # base64-encoded JPEG


class FaceVerifyResponse(BaseModel):
    verified: bool
    confidence: float
    student_id: str
    message: str


def _cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    dot = float(np.dot(a, b))
    norm_a = float(np.linalg.norm(a))
    norm_b = float(np.linalg.norm(b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def _sync_face_enroll(student_id: str, image_b64: str) -> FaceEnrollResponse:
    """Extract ArcFace embedding and store on disk."""
    try:
        DeepFace = _get_deepface()
        image_data = base64.b64decode(image_b64)
        image = Image.open(io.BytesIO(image_data))
        frame = np.array(image)
        frame_bgr = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)

        embeddings = DeepFace.represent(
            img_path=frame_bgr,
            model_name="ArcFace",
            enforce_detection=True,
            detector_backend="mediapipe",
        )

        if not embeddings:
            return FaceEnrollResponse(
                success=False, message="No face detected in image", student_id=student_id
            )

        embedding = np.array(embeddings[0]["embedding"], dtype=np.float32)
        embedding_path = os.path.join(FACE_EMBEDDINGS_DIR, f"{student_id}.npy")
        np.save(embedding_path, embedding)

        return FaceEnrollResponse(
            success=True, message="Face enrolled successfully", student_id=student_id
        )
    except Exception as e:
        return FaceEnrollResponse(
            success=False, message=f"Enrollment failed: {str(e)}", student_id=student_id
        )


@app.post("/api/face/enroll", response_model=FaceEnrollResponse)
async def face_enroll(request: FaceEnrollRequest):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, partial(_sync_face_enroll, request.student_id, request.image)
    )


def _sync_face_verify(student_id: str, image_b64: str) -> FaceVerifyResponse:
    """Verify a face against the stored enrollment embedding."""
    embedding_path = os.path.join(FACE_EMBEDDINGS_DIR, f"{student_id}.npy")

    if not os.path.exists(embedding_path):
        return FaceVerifyResponse(
            verified=False, confidence=0.0, student_id=student_id,
            message="No enrolled face found. Please enroll first.",
        )

    stored_embedding = np.load(embedding_path)

    try:
        DeepFace = _get_deepface()
        image_data = base64.b64decode(image_b64)
        image = Image.open(io.BytesIO(image_data))
        frame = np.array(image)
        frame_bgr = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)

        embeddings = DeepFace.represent(
            img_path=frame_bgr,
            model_name="ArcFace",
            enforce_detection=True,
            detector_backend="mediapipe",
        )

        if not embeddings:
            return FaceVerifyResponse(
                verified=False, confidence=0.0, student_id=student_id,
                message="No face detected in verification image",
            )

        current_embedding = np.array(embeddings[0]["embedding"], dtype=np.float32)
        similarity = _cosine_similarity(stored_embedding, current_embedding)

        verified = similarity >= ARCFACE_THRESHOLD
        return FaceVerifyResponse(
            verified=verified,
            confidence=round(similarity, 4),
            student_id=student_id,
            message="Identity verified" if verified else f"Identity mismatch (confidence: {similarity:.2%})",
        )
    except Exception as e:
        return FaceVerifyResponse(
            verified=False, confidence=0.0, student_id=student_id,
            message=f"Verification failed: {str(e)}",
        )


@app.post("/api/face/verify", response_model=FaceVerifyResponse)
async def face_verify(request: FaceVerifyRequest):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None, partial(_sync_face_verify, request.student_id, request.image)
    )


@app.get("/api/face/check/{student_id}")
async def check_face_enrollment(student_id: str):
    embedding_path = os.path.join(FACE_EMBEDDINGS_DIR, f"{student_id}.npy")
    return {"enrolled": os.path.exists(embedding_path), "student_id": student_id}


# ── Liveness Detection ─────────────────────────────────────────────

class LivenessFrame(BaseModel):
    image: str
    timestamp: float  # milliseconds since challenge start

class LivenessRequest(BaseModel):
    student_id: str
    challenge_type: str  # "BLINK" or "HEAD_TURN"
    challenge_param: str  # "3" for blink count, "LEFT" or "RIGHT" for head turn
    frames: List[LivenessFrame]

class LivenessResponse(BaseModel):
    passed: bool
    confidence: float
    message: str
    details: Optional[Dict] = None


def _compute_ear(landmarks) -> float:
    """Compute average Eye Aspect Ratio from landmarks."""
    def ear(eye_indices):
        pts = [(landmarks[i].x, landmarks[i].y) for i in eye_indices]
        v1 = math.dist(pts[1], pts[5])
        v2 = math.dist(pts[2], pts[4])
        h = math.dist(pts[0], pts[3])
        return (v1 + v2) / (2.0 * h) if h > 0 else 0

    left_eye = [33, 160, 158, 133, 153, 144]
    right_eye = [362, 385, 387, 263, 373, 380]
    return (ear(left_eye) + ear(right_eye)) / 2.0


def _sync_liveness_check(student_id: str, challenge_type: str, challenge_param: str, frames: List[dict]) -> dict:
    """Analyze a sequence of frames for liveness challenge response."""
    if len(frames) < 3:
        return {"passed": False, "confidence": 0.0, "message": "Too few frames provided", "details": None}

    ear_values = []
    yaw_values = []
    valid_frames = 0

    for frame_data in frames:
        try:
            image_data = base64.b64decode(frame_data["image"])
            image = Image.open(io.BytesIO(image_data))
            frame = np.array(image)
            if len(frame.shape) == 3 and frame.shape[2] == 4:
                frame = cv2.cvtColor(frame, cv2.COLOR_RGBA2RGB)
            elif len(frame.shape) == 3 and frame.shape[2] == 3:
                # Check if BGR (from cv2) or RGB (from PIL) — PIL gives RGB
                pass  # Already RGB from PIL
            mp_image = mp.Image(image_format=mp.ImageFormat.SRGB, data=frame)
            landmarker_result = face_landmarker.detect(mp_image)

            if not landmarker_result.face_landmarks:
                continue

            landmarks = landmarker_result.face_landmarks[0]
            valid_frames += 1

            # Compute EAR for blink detection
            avg_ear = _compute_ear(landmarks)
            ear_values.append(avg_ear)

            # Compute yaw for head turn detection
            h, w = frame.shape[:2]
            pose = estimate_head_pose(landmarks, w, h)
            yaw_values.append(pose["yaw"])

        except Exception:
            continue

    if valid_frames < 3:
        return {"passed": False, "confidence": 0.1, "message": "Could not detect face in enough frames", "details": {"valid_frames": valid_frames}}

    if challenge_type == "BLINK":
        # Count blink events: EAR drops below 0.18 then rises above 0.22
        required_blinks = int(challenge_param) if challenge_param.isdigit() else 3
        blink_count = 0
        in_blink = False
        for ear_val in ear_values:
            if not in_blink and ear_val < 0.18:
                in_blink = True
            elif in_blink and ear_val > 0.22:
                blink_count += 1
                in_blink = False

        passed = blink_count >= required_blinks
        confidence = min(1.0, blink_count / required_blinks) if required_blinks > 0 else 0.0
        return {
            "passed": passed,
            "confidence": round(confidence, 2),
            "message": f"Detected {blink_count}/{required_blinks} blinks" if passed else f"Only detected {blink_count}/{required_blinks} blinks",
            "details": {"blink_count": blink_count, "required": required_blinks, "valid_frames": valid_frames},
        }

    elif challenge_type == "HEAD_TURN":
        # Check if yaw exceeded ±20° in the expected direction
        direction = challenge_param.upper()
        if direction == "LEFT":
            max_yaw = max(yaw_values) if yaw_values else 0
            passed = max_yaw > 20.0
            confidence = min(1.0, max_yaw / 20.0) if max_yaw > 0 else 0.0
        elif direction == "RIGHT":
            min_yaw = min(yaw_values) if yaw_values else 0
            passed = min_yaw < -20.0
            confidence = min(1.0, abs(min_yaw) / 20.0) if min_yaw < 0 else 0.0
        else:
            return {"passed": False, "confidence": 0.0, "message": f"Unknown direction: {direction}", "details": None}

        return {
            "passed": passed,
            "confidence": round(confidence, 2),
            "message": f"Head turn {direction.lower()} detected" if passed else f"Head turn {direction.lower()} not detected",
            "details": {"direction": direction, "yaw_range": [round(min(yaw_values), 1), round(max(yaw_values), 1)], "valid_frames": valid_frames},
        }

    return {"passed": False, "confidence": 0.0, "message": f"Unknown challenge type: {challenge_type}", "details": None}


@app.post("/api/face/liveness-check", response_model=LivenessResponse)
async def liveness_check(request: LivenessRequest):
    loop = asyncio.get_event_loop()
    frames_data = [{"image": f.image, "timestamp": f.timestamp} for f in request.frames]
    result = await loop.run_in_executor(
        None,
        partial(_sync_liveness_check, request.student_id, request.challenge_type, request.challenge_param, frames_data),
    )
    return LivenessResponse(**result)


if __name__ == "__main__":
    import uvicorn
    import sys

    port = 8001
    if len(sys.argv) > 1:
        try:
            port = int(sys.argv[1])
        except ValueError:
            pass

    uvicorn.run(app, host="0.0.0.0", port=port)
